import os
import json
from pathlib import Path
import requests
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import chromadb

###########################################################
# Configuration
###########################################################
ROOT_FOLDER = Path.home() / "Desktop" / "CourseMaterials"

# NVIDIA NIM API (OpenAI-compatible embeddings endpoint)
NVIDIA_API_KEY = os.environ["NVIDIA_API_KEY"]  # set this in your shell/env
NVIDIA_API_URL = "https://integrate.api.nvidia.com/v1/embeddings"
EMBEDDING_MODEL = "nvidia/nv-embedqa-e5-v5"

CHUNK_SIZE = 1000
OVERLAP = 200
CHROMA_DB = "./course_kb"

###########################################################
# Chroma DB
###########################################################
client = chromadb.PersistentClient(path=CHROMA_DB)
collection = client.get_or_create_collection(name="course_materials")

###########################################################
# Readers
###########################################################
def read_pdf(path):
    reader = PdfReader(path)
    text = ""
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text


def read_docx(path):
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs)


def read_pptx(path):
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)


def read_txt(path):
    with open(path, "r", encoding="utf-8") as f:
        return f.read()


###########################################################
# Chunking
###########################################################
def chunk_text(text):
    chunks = []
    start = 0
    while start < len(text):
        end = start + CHUNK_SIZE
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start += CHUNK_SIZE - OVERLAP
    return chunks


###########################################################
# NVIDIA NIM Embeddings
###########################################################
def get_embedding(text, input_type="passage"):
    """
    input_type must be 'passage' when embedding documents for storage,
    and 'query' when embedding a user's search query at retrieval time.
    NVIDIA's E5 embedding models require this to be set correctly.
    """
    headers = {
        "Authorization": f"Bearer {NVIDIA_API_KEY}",
        "Content-Type": "application/json",
        "Accept": "application/json",
    }
    payload = {
        "input": [text],
        "model": EMBEDDING_MODEL,
        "input_type": input_type,
        "encoding_format": "float",
    }

    response = requests.post(NVIDIA_API_URL, headers=headers, json=payload, timeout=30)
    response.raise_for_status()
    data = response.json()
    return data["data"][0]["embedding"]


###########################################################
# Main
###########################################################
supported = {
    ".pdf": read_pdf,
    ".docx": read_docx,
    ".pptx": read_pptx,
    ".txt": read_txt,
}

for subject in ROOT_FOLDER.iterdir():
    if not subject.is_dir():
        continue

    subject_name = subject.name
    print(f"\nProcessing {subject_name}")

    for file in subject.rglob("*"):
        if file.suffix.lower() not in supported:
            continue

        print(file.name)
        extractor = supported[file.suffix.lower()]

        try:
            text = extractor(file)
        except Exception as e:
            print(f"  ! failed to read {file.name}: {e}")
            continue

        chunks = chunk_text(text)

        # Unique, stable ID per subject+file+chunk. Using upsert (not add)
        # means re-running the script safely updates existing chunks
        # instead of erroring or colliding with chunks from other subjects.
        safe_filename = file.stem.replace(" ", "_")
        doc_id_base = f"{subject_name}_{safe_filename}"

        for idx, chunk in enumerate(chunks):
            try:
                embedding = get_embedding(chunk, input_type="passage")
            except Exception as e:
                print(f"  ! embedding failed for chunk {idx} of {file.name}: {e}")
                continue

            collection.upsert(
                ids=[f"{doc_id_base}_{idx}"],
                embeddings=[embedding],
                documents=[chunk],
                metadatas=[{
                    "subject": subject_name,
                    "filename": file.name,
                    "path": str(file),
                }],
            )

print("\nKnowledge Base Created Successfully.")